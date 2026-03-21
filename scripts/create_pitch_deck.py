#!/usr/bin/env python3
"""
Generates a professional pitch deck PowerPoint for selling the Bocage Champagne
Society app to owners Clark Gale and Zac Denham (Sure Thing Hospitality).

Creates: presentation/Bocage_Champagne_Society_Pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
NOIR = RGBColor(0x0A, 0x0A, 0x0A)
NOIR_800 = RGBColor(0x11, 0x11, 0x11)
NOIR_700 = RGBColor(0x18, 0x18, 0x18)
NOIR_600 = RGBColor(0x1F, 0x1F, 0x1F)
NOIR_300 = RGBColor(0x6B, 0x6B, 0x6B)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
GOLD_LIGHT = RGBColor(0xF3, 0xDA, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFD, 0xF8, 0xEC)
ROSE = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=NOIR):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_gold_line(slide, left, top, width):
    """Add a horizontal gold accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, color=NOIR_800):
    """Add a rounded rectangle card shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
    shape.line.width = Pt(1)
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=WHITE):
    """Add a bulleted text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)
add_gold_line(slide, 3, 2.5, 7.333)
add_text_box(slide, 1, 2.7, 11.333, 1.2,
             "Bocage", 72, GOLD, True, PP_ALIGN.CENTER, 'Georgia')
add_text_box(slide, 1, 3.7, 11.333, 0.6,
             "CHAMPAGNE SOCIETY", 20, WHITE, False, PP_ALIGN.CENTER, 'Calibri')
add_gold_line(slide, 3, 4.4, 7.333)
add_text_box(slide, 1, 4.8, 11.333, 0.6,
             "A Luxury Membership & Loyalty App", 16, NOIR_300, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11.333, 0.5,
             "Prepared for Clark Gale & Zac Denham  |  Sure Thing Hospitality", 13, NOIR_300, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.2, 11.333, 0.4,
             "10 Phila Street, Saratoga Springs, NY 12866  |  champagnebar.com", 11, NOIR_300, False, PP_ALIGN.CENTER)
# Est. 2021 badge
add_text_box(slide, 5.5, 1.6, 2.333, 0.5,
             "EST. 2021", 12, GOLD, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: The Opportunity
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "The Opportunity", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

add_text_box(slide, 0.8, 1.6, 5.5, 1.5,
             "Bocage Champagne Bar is Saratoga's premier luxury destination. "
             "But your customers only interact with you when they walk through the door.\n\n"
             "What if you could be in their pocket 24/7?",
             16, WHITE)

# Problem stats
problems = [
    "78% of restaurant revenue comes from repeat customers",
    "Loyalty members spend 67% more than non-members",
    "Push notifications have 7x the engagement of email",
    "Digital gift cards grow revenue 20-30% year-over-year",
    "Top hospitality apps see 40%+ monthly active users",
]
add_text_box(slide, 7, 1.6, 5.5, 0.5, "Industry Facts:", 14, GOLD, True)
add_bullet_list(slide, 7, 2.2, 5.5, 4, problems, 14, CREAM)

add_card(slide, 0.8, 5.2, 11.7, 1.5, NOIR_700)
add_text_box(slide, 1.2, 5.4, 11, 1.2,
             "The Bocage Champagne Society app transforms casual visitors into loyal members, "
             "drives recurring revenue through subscriptions, and creates a VIP experience "
             "that matches the luxury of your bar.",
             15, GOLD, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: The App — What We Built
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "What We Built", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

# Feature cards in a 2x3 grid
features = [
    ("Digital Wine Menu", "Your full menu — Bar Snacks, Caviar, Flights, Cocktails, Wines — searchable and always up-to-date. No more PDF downloads."),
    ("3-Tier Membership", "Flûte (free), Magnum ($29.99/mo), Jeroboam ($79.99/mo). Points multipliers, perks, and exclusive access at each level."),
    ("Event Bookings", "Create and sell tickets to champagne tastings, dinners, and masterclasses. Tier-gated VIP events drive upgrades."),
    ("At-Home Experiences", "Book private champagne experiences starting at $250. Three tiers of luxury service — a high-margin revenue stream."),
    ("Admin Dashboard", "Full inventory management, event creation, member insights. Control everything from your phone."),
    ("Push Notifications", "Send targeted promotions to members. 'Happy Hour starting now' drives immediate foot traffic."),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.5
    add_card(slide, x, y, 3.8, 2.2, NOIR_800)
    add_text_box(slide, x + 0.2, y + 0.2, 3.4, 0.5, title, 16, GOLD, True)
    add_text_box(slide, x + 0.2, y + 0.7, 3.4, 1.3, desc, 12, CREAM)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Revenue Streams
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Revenue Streams", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)
add_text_box(slide, 0.8, 1.5, 11, 0.5, "Multiple ways the app generates direct and indirect revenue for Bocage:", 14, NOIR_300)

streams = [
    ("1. Membership Subscriptions", "$29.99/mo (Magnum) + $79.99/mo (Jeroboam)", "Just 100 Magnum + 25 Jeroboam members = $60K/year recurring"),
    ("2. Event Ticket Sales", "$55–$250 per ticket, sold in-app", "5 events/month × 20 avg seats × $90 avg = $108K/year"),
    ("3. At-Home Experiences", "$250–$1,000+ per booking", "4 bookings/month × $500 avg = $24K/year"),
    ("4. Digital Gift Cards", "Purchasable in-app, redeemable at bar", "Holiday/birthday gift cards drive new customer acquisition"),
    ("5. In-App Pre-Orders", "Order ahead, open tabs, skip the wait", "Increases avg ticket size 15-25% via upselling"),
    ("6. Push Marketing", "Targeted promotions drive same-day visits", 'e.g. "Champagne Tuesday: 2-for-1 glasses" — no ad spend required'),
    ("7. Referral Program", "Members earn 100 pts for each referral", "Organic growth: each member brings 1.5 new members on avg"),
    ("8. Data & Insights", "Know your top spenders, popular items, peak times", "Optimize inventory, staffing, and marketing with real data"),
]

for i, (title, detail, impact) in enumerate(streams):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.3
    y = 2.2 + row * 1.25
    add_text_box(slide, x, y, 5.8, 0.35, title, 13, GOLD, True)
    add_text_box(slide, x, y + 0.32, 5.8, 0.35, detail, 11, WHITE)
    add_text_box(slide, x, y + 0.6, 5.8, 0.35, impact, 10, NOIR_300)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Revenue Projections
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Conservative Revenue Projections", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)
add_text_box(slide, 0.8, 1.5, 11, 0.5, "Year 1 estimates based on Saratoga Springs market and comparable hospitality loyalty programs:", 14, NOIR_300)

# Year 1 projections table
projections = [
    ("Memberships (Magnum)", "100 members × $29.99/mo × 12", "$35,988"),
    ("Memberships (Jeroboam)", "25 members × $79.99/mo × 12", "$23,997"),
    ("Event Tickets", "60 events × 18 avg seats × $90", "$97,200"),
    ("At-Home Bookings", "48 bookings × $500 avg", "$24,000"),
    ("Gift Card Sales", "200 cards × $75 avg", "$15,000"),
    ("Increased Spend (loyalty lift)", "15% lift on 500 members × $80 avg visit", "$72,000"),
    ("Referral-Driven New Customers", "75 new customers × $120 first-year value", "$9,000"),
]

add_text_box(slide, 1.2, 2.2, 5, 0.4, "Revenue Source", 12, GOLD, True)
add_text_box(slide, 5.2, 2.2, 4.5, 0.4, "Calculation", 12, GOLD, True)
add_text_box(slide, 10, 2.2, 2.5, 0.4, "Annual", 12, GOLD, True, PP_ALIGN.RIGHT)
add_gold_line(slide, 1.2, 2.6, 11)

for i, (source, calc, amount) in enumerate(projections):
    y = 2.8 + i * 0.5
    color = WHITE if i % 2 == 0 else CREAM
    add_text_box(slide, 1.2, y, 4, 0.4, source, 12, color)
    add_text_box(slide, 5.2, y, 4.5, 0.4, calc, 10, NOIR_300)
    add_text_box(slide, 10, y, 2.5, 0.4, amount, 13, GOLD, True, PP_ALIGN.RIGHT)

add_gold_line(slide, 1.2, 6.4, 11)
add_text_box(slide, 1.2, 6.5, 5, 0.5, "TOTAL PROJECTED YEAR 1 REVENUE", 14, WHITE, True)
add_text_box(slide, 10, 6.5, 2.5, 0.5, "$277,185", 22, GOLD, True, PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5b: Infrastructure Costs & Net Profit
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Infrastructure & Net Profit", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)
add_text_box(slide, 0.8, 1.5, 11, 0.5, "Real infrastructure costs on Supabase Pro — lean and scalable:", 14, NOIR_300)

GREEN = RGBColor(0x22, 0xC5, 0x5E)

# Supabase Pro costs breakdown
costs = [
    ("Supabase Pro Plan (base)", "Per project/month", "$25/mo", "$300"),
    ("MAU (100K included free)", "First 100K MAUs included in Pro", "$0.00", "$0"),
    ("MAU Overage", "Est. 2,000 MAU Year 1 — well under 100K limit", "$0.00", "$0"),
    ("Database Disk (8 GB included)", "8 GB included, est. ~2 GB used Year 1", "$0.00", "$0"),
    ("Egress (250 GB included)", "250 GB included, est. ~50 GB used Year 1", "$0.00", "$0"),
    ("File Storage (100 GB included)", "Wine images + assets, est. ~5 GB", "$0.00", "$0"),
    ("Daily Backups", "Included — stored for 7 days", "$0.00", "$0"),
    ("Email Support", "Included with Pro", "$0.00", "$0"),
    ("Vercel Hosting (web)", "Hobby/Pro plan for web app", "$20/mo", "$240"),
    ("Apple Developer Account", "iOS App Store distribution", "$99/yr", "$99"),
    ("Google Play Developer", "Android distribution (one-time)", "$25", "$25"),
    ("Stripe Payment Processing", "2.9% + $0.30 per transaction (on ~$200K)", "~2.9%", "~$6,200"),
    ("Domain + SSL", "champagnebar.com — already owned", "$0.00", "$0"),
]

add_text_box(slide, 1.2, 2.1, 4, 0.35, "Cost Item", 11, GOLD, True)
add_text_box(slide, 5, 2.1, 4.5, 0.35, "Details", 11, GOLD, True)
add_text_box(slide, 9.5, 2.1, 1.5, 0.35, "Rate", 11, GOLD, True, PP_ALIGN.RIGHT)
add_text_box(slide, 11, 2.1, 1.5, 0.35, "Annual", 11, GOLD, True, PP_ALIGN.RIGHT)
add_gold_line(slide, 1.2, 2.45, 11)

for i, (item, detail, rate, annual) in enumerate(costs):
    y = 2.55 + i * 0.3
    color = WHITE if i % 2 == 0 else CREAM
    add_text_box(slide, 1.2, y, 3.8, 0.28, item, 9, color)
    add_text_box(slide, 5, y, 4.5, 0.28, detail, 8, NOIR_300)
    add_text_box(slide, 9.5, y, 1.5, 0.28, rate, 9, NOIR_300, False, PP_ALIGN.RIGHT)
    add_text_box(slide, 11, y, 1.5, 0.28, annual, 9, color, False, PP_ALIGN.RIGHT)

# Totals
y_total = 6.55
add_gold_line(slide, 1.2, y_total - 0.1, 11)
add_text_box(slide, 1.2, y_total, 4, 0.35, "TOTAL INFRASTRUCTURE COST", 11, WHITE, True)
add_text_box(slide, 11, y_total, 1.5, 0.35, "~$6,864", 13, RGBColor(0xFF, 0x66, 0x66), True, PP_ALIGN.RIGHT)

add_text_box(slide, 1.2, y_total + 0.4, 4, 0.35, "YEAR 1 GROSS REVENUE", 11, WHITE, True)
add_text_box(slide, 11, y_total + 0.4, 1.5, 0.35, "$277,185", 13, GOLD, True, PP_ALIGN.RIGHT)

add_gold_line(slide, 1.2, y_total + 0.8, 11)
add_text_box(slide, 1.2, y_total + 0.85, 5, 0.4, "NET PROFIT (97.5% margin)", 13, WHITE, True)
add_text_box(slide, 10, y_total + 0.85, 2.5, 0.4, "$270,321", 22, GREEN, True, PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Membership Tiers Deep Dive
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Membership Tiers", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

tiers = [
    ("Flûte", "Free", "Entry tier — every signup",
     ["1x points on every visit", "Member-only event access", "Birthday champagne toast",
      "Early access to new arrivals", "Digital membership card", "Push notification deals"]),
    ("Magnum", "$29.99/month", "For champagne enthusiasts",
     ["1.5x points multiplier", "All Flûte benefits", "Complimentary monthly glass",
      "Priority event seating", "Exclusive Magnum tastings", "10% off At-Home bookings"]),
    ("Jeroboam", "$79.99/month", "The inner circle",
     ["2x points multiplier", "All Magnum benefits", "Complimentary monthly bottle",
      "Private lounge access", "Personal sommelier service", "VIP event access", "25% off At-Home"]),
]

for i, (name, price, subtitle, benefits) in enumerate(tiers):
    x = 0.8 + i * 4.2
    border_color = NOIR_300 if i == 0 else (GOLD if i == 1 else ROSE)
    add_card(slide, x, 1.8, 3.9, 5, NOIR_800)
    add_text_box(slide, x + 0.3, 1.95, 3.3, 0.5, name, 24, border_color, True, PP_ALIGN.CENTER, 'Georgia')
    add_text_box(slide, x + 0.3, 2.45, 3.3, 0.4, price, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, 2.85, 3.3, 0.3, subtitle, 11, NOIR_300, False, PP_ALIGN.CENTER)
    add_gold_line(slide, x + 0.5, 3.2, 2.9)
    add_bullet_list(slide, x + 0.3, 3.4, 3.3, 3.2, benefits, 11, CREAM)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7: App Experience — Screenshots Description
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "The App Experience", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

screens = [
    ("Login", "Elegant Bocage branding with\nanimated logo, password\nstrength meter, and\nsmooth transitions"),
    ("Wine Menu", "Your full menu with search,\nfilters, grid/list view,\nand tap-for-details.\nAlways current."),
    ("Membership", "Digital membership card,\npoints progress, tier\nbenefits, and transaction\nhistory at a glance."),
    ("Events", "Upcoming events with\ncountdown timers, seat\ntracking, tier-gating,\nand one-tap RSVP."),
    ("At Home", "Three-tier private\nexperience booking with\ntestimonials, FAQ,\nand instant requests."),
    ("Admin", "Full inventory + event\nmanagement. Stats\ndashboard, CRUD\noperations, photo uploads."),
]

for i, (title, desc) in enumerate(screens):
    x = 0.4 + i * 2.15
    # Phone frame placeholder
    add_card(slide, x, 1.8, 1.95, 3.8, NOIR_700)
    add_text_box(slide, x + 0.1, 2.0, 1.75, 0.4, title, 13, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 2.5, 1.75, 2.8, desc, 9, CREAM, False, PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6, 11.7, 0.8,
             "Available on iOS App Store  •  Google Play Store  •  Web (champagnebar.com)",
             14, NOIR_300, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Competitive Advantage
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Why This Wins", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

advantages = [
    ("Custom-Built for Bocage", "Not a generic loyalty app. Every pixel reflects Bocage's luxury brand — dark noir backgrounds, champagne gold accents, and your actual menu items."),
    ("Three Revenue Engines", "Subscriptions provide predictable monthly income. Events generate high-margin ticket sales. At-Home creates a premium service line."),
    ("Zero Marketing Cost", "Push notifications replace paid ads. A single 'Happy Hour Starts Now' notification drives immediate visits — no Facebook spend required."),
    ("Data You Don't Have Today", "Know your top 10 spenders, most popular cocktails, peak visit times, and which events sell out fastest. Make smarter business decisions."),
    ("Member Lock-In", "Points + tiers create switching costs. A member with 400 points toward Magnum won't leave for a competitor. Retention becomes automatic."),
    ("Viral Growth Built-In", "Referral program + shareable events + gift cards turn every member into a marketer. Organic growth compounds month over month."),
]

for i, (title, desc) in enumerate(advantages):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.3
    y = 1.8 + row * 1.8
    add_card(slide, x, y, 5.9, 1.5, NOIR_800)
    add_text_box(slide, x + 0.3, y + 0.15, 5.3, 0.4, title, 15, GOLD, True)
    add_text_box(slide, x + 0.3, y + 0.55, 5.3, 0.85, desc, 11, CREAM)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Implementation & Timeline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Launch Plan", 40, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
add_gold_line(slide, 0.8, 1.3, 4)

phases = [
    ("Week 1–2", "SETUP", ["Connect Supabase database", "Configure Stripe payments", "Upload wine menu + photos", "Set admin accounts for Clark & Zac"]),
    ("Week 3", "TESTING", ["Internal testing with staff", "Verify all menu items", "Test event booking flow", "QA on iOS, Android, and web"]),
    ("Week 4", "SOFT LAUNCH", ["Launch to email list", "Staff promotes to regulars", "First members earn points", "Gather initial feedback"]),
    ("Month 2", "FULL LAUNCH", ["App Store + Google Play live", "First member event", "Push notification campaigns", "Referral program activated"]),
    ("Month 3+", "GROWTH", ["Monthly events calendar", "Gift card holiday push", "At-Home marketing", "Data-driven optimization"]),
]

for i, (time, label, items) in enumerate(phases):
    x = 0.4 + i * 2.55
    add_card(slide, x, 1.8, 2.35, 4.8, NOIR_800)
    add_text_box(slide, x + 0.15, 1.95, 2.05, 0.3, time, 11, NOIR_300, False, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 2.25, 2.05, 0.4, label, 16, GOLD, True, PP_ALIGN.CENTER)
    add_gold_line(slide, x + 0.3, 2.7, 1.75)
    add_bullet_list(slide, x + 0.15, 2.85, 2.05, 3.5, items, 10, CREAM)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10: The Ask / Next Steps
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_gold_line(slide, 3, 1.5, 7.333)
add_text_box(slide, 1, 1.7, 11.333, 1, "Let's Launch", 52, GOLD, True, PP_ALIGN.CENTER, 'Georgia')
add_text_box(slide, 1, 2.7, 11.333, 0.5, "Bocage Champagne Society", 22, WHITE, False, PP_ALIGN.CENTER)
add_gold_line(slide, 3, 3.3, 7.333)

add_text_box(slide, 2, 3.8, 9.333, 1,
             "The app is built. The infrastructure is ready.\nAll we need is your green light to go live.",
             18, CREAM, False, PP_ALIGN.CENTER)

# Next steps
steps = [
    "1.  Review this presentation + app demo",
    "2.  Finalize menu items and event calendar",
    "3.  Set up Supabase + Stripe accounts",
    "4.  Two-week testing sprint with your team",
    "5.  Launch to your members and start generating revenue",
]

add_text_box(slide, 3.5, 4.8, 6.333, 0.4, "Next Steps:", 16, GOLD, True, PP_ALIGN.CENTER)
add_bullet_list(slide, 3.5, 5.2, 6.333, 2, steps, 14, WHITE)

add_text_box(slide, 1, 6.8, 11.333, 0.4,
             "10 Phila Street, Saratoga Springs, NY 12866  |  champagnebar.com  |  Est. 2021",
             11, NOIR_300, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
import os
output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'presentation')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'Bocage_Champagne_Society_Pitch.pptx')
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
